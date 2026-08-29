"""work-kb - Office parsing sidecar.

JSON-RPC over stdio (one request/response per line).
Request:  {"id":1,"method":"parse","params":{"path":"...","kind":"docx|xlsx|pptx|doc|xls|ppt|rtf|wps|et|dps"}}
Response: {"id":1,"ok":true,"result":{ParseResult}} or {"id":1,"ok":false,"error":"..."}

Dependencies: python-docx / openpyxl / python-pptx / xlrd / striprtf / olefile (see requirements.txt).
Libraries are lazily imported per handler, so a missing lib does not affect other types.
"""

import json
import os
import sys


# ── OOXML handlers ──────────────────────────────────────────────

def parse_docx(path):
    """DOCX: python-docx, headings + paragraphs + tables."""
    from docx import Document

    doc = Document(path)
    title = (doc.core_properties.title or
             os.path.splitext(os.path.basename(path))[0] or "Word")
    sections = []
    cur = {"heading": "", "level": 0, "body": "", "page": None}
    started = False

    for p in doc.paragraphs:
        style = (p.style.name or "") if p.style else ""
        text = p.text or ""
        if style.startswith("Heading") or style == "Title":
            if started:
                sections.append(cur)
            if style == "Title":
                lvl = 1
            else:
                try:
                    lvl = int(style.replace("Heading", "").strip())
                except ValueError:
                    lvl = 1
            cur = {"heading": text, "level": lvl, "body": "", "page": None}
            started = True
        else:
            if text.strip():
                cur["body"] += text + "\n"
                started = True

    # Extract tables as Markdown table format
    for table in doc.tables:
        if started:
            sections.append(cur)
            cur = {"heading": "", "level": 0, "body": "", "page": None}
            started = False
        table_lines = []
        for row in table.rows:
            cells = [cell.text.strip().replace("\n", " ") for cell in row.cells]
            table_lines.append("| " + " | ".join(cells) + " |")
        if table_lines:
            cur["body"] = "\n".join(table_lines) + "\n"
            started = True

    if started:
        sections.append(cur)
    if not sections:
        sections = [{"heading": title, "level": 0, "body": "", "page": None}]
    return {"sourcePath": path, "docTitle": title, "sections": sections}


def parse_xlsx(path):
    """XLSX: openpyxl, one section per sheet, TSV rows."""
    from openpyxl import load_workbook

    wb = load_workbook(path, data_only=True, read_only=True)
    title = os.path.splitext(os.path.basename(path))[0] or "Excel"
    sections = []
    for ws in wb.worksheets:
        lines = []
        for row in ws.iter_rows(values_only=True):
            cells = ["" if c is None else str(c) for c in row]
            if any(c.strip() for c in cells):
                lines.append("\t".join(cells))
        sections.append({
            "heading": ws.title or "Sheet",
            "level": 0,
            "body": "\n".join(lines),
            "page": None,
        })
    return {"sourcePath": path, "docTitle": title, "sections": sections}


def parse_pptx(path):
    """PPTX: python-pptx, one section per slide, includes speaker notes."""
    from pptx import Presentation

    prs = Presentation(path)
    title = (prs.core_properties.title or
             os.path.splitext(os.path.basename(path))[0] or "PPT")
    sections = []
    for i, slide in enumerate(prs.slides, 1):
        title_shape = slide.shapes.title
        heading = ""
        if title_shape is not None and title_shape.has_text_frame:
            heading = title_shape.text_frame.text or ""
        body_parts = []
        for shape in slide.shapes:
            if shape is title_shape:
                continue
            if shape.has_text_frame:
                txt = shape.text_frame.text
                if txt.strip():
                    body_parts.append(txt)
        # Extract speaker notes (P2 improvement)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame
            if notes and notes.text.strip():
                body_parts.append("[notes] " + notes.text.strip())
        if not heading:
            heading = f"Slide {i}"
        sections.append({
            "heading": heading,
            "level": 1,
            "body": "\n".join(body_parts),
            "page": i,
        })
    return {"sourcePath": path, "docTitle": title, "sections": sections}


# ── Legacy Office handlers ──────────────────────────────────────

def parse_doc(path):
    """DOC: legacy Word. Try antiword, fall back to olefile text extraction."""
    title = os.path.splitext(os.path.basename(path))[0] or "Word"

    # Try antiword
    try:
        import subprocess
        result = subprocess.run(
            ["antiword", path],
            capture_output=True, text=True, timeout=30,
        )
        if result.returncode == 0 and result.stdout.strip():
            return {"sourcePath": path, "docTitle": title, "sections": [
                {"heading": "", "level": 0, "body": result.stdout.strip(), "page": None}
            ]}
    except Exception:
        pass

    # Fallback: olefile + UTF-16 text extraction from WordDocument stream
    try:
        import olefile
        ole = olefile.OleFileIO(path)
        if ole.exists("WordDocument"):
            data = ole.openstream("WordDocument").read()
            ole.close()
            text = data.decode("utf-16-le", errors="ignore")
            lines = [l.strip() for l in text.split("\n")
                     if l.strip() and len(l.strip()) > 2]
            if lines:
                return {"sourcePath": path, "docTitle": title, "sections": [
                    {"heading": "", "level": 0,
                     "body": "\n".join(lines[:500]), "page": None}
                ]}
        ole.close()
    except Exception:
        pass

    # Last resort: raw binary text extraction
    with open(path, "rb") as f:
        data = f.read()
    text = data.decode("utf-16-le", errors="ignore")
    lines = [l.strip() for l in text.split("\n")
             if l.strip() and len(l.strip()) > 5]
    if lines:
        return {"sourcePath": path, "docTitle": title, "sections": [
            {"heading": "", "level": 0,
             "body": "\n".join(lines[:200]), "page": None}
        ]}

    raise RuntimeError("Cannot parse DOC. Please save as DOCX and retry.")


def parse_xls(path):
    """XLS: legacy Excel. Uses xlrd."""
    import xlrd

    title = os.path.splitext(os.path.basename(path))[0] or "Excel"
    wb = xlrd.open_workbook(path)
    sections = []
    for ws in wb.sheets():
        lines = []
        for row_idx in range(ws.nrows):
            cells = [str(ws.cell_value(row_idx, c)) for c in range(ws.ncols)]
            if any(c.strip() for c in cells):
                lines.append("\t".join(cells))
        sections.append({
            "heading": ws.name or "Sheet",
            "level": 0,
            "body": "\n".join(lines),
            "page": None,
        })
    return {"sourcePath": path, "docTitle": title, "sections": sections}


def parse_ppt(path):
    """PPT: legacy PowerPoint. Try LibreOffice conversion, fall back to olefile."""
    title = os.path.splitext(os.path.basename(path))[0] or "PPT"

    # Try LibreOffice headless conversion
    try:
        import subprocess
        import tempfile
        with tempfile.TemporaryDirectory() as tmpdir:
            result = subprocess.run(
                ["libreoffice", "--headless", "--convert-to", "pptx",
                 "--outdir", tmpdir, path],
                capture_output=True, timeout=60,
            )
            if result.returncode == 0:
                basename = os.path.splitext(os.path.basename(path))[0]
                pptx_path = os.path.join(tmpdir, basename + ".pptx")
                if os.path.exists(pptx_path):
                    result = parse_pptx(pptx_path)
                    result["sourcePath"] = path
                    return result
    except Exception:
        pass

    # Fallback: olefile text extraction from all streams
    try:
        import olefile
        ole = olefile.OleFileIO(path)
        text_parts = []
        for stream in ole.listdir():
            try:
                data = ole.openstream(stream).read()
                decoded = data.decode("utf-16-le", errors="ignore")
                lines = [l.strip() for l in decoded.split("\n")
                         if l.strip() and len(l.strip()) > 3]
                if lines:
                    text_parts.extend(lines)
            except Exception:
                continue
        ole.close()
        if text_parts:
            return {"sourcePath": path, "docTitle": title, "sections": [
                {"heading": "", "level": 0,
                 "body": "\n".join(text_parts[:500]), "page": None}
            ]}
    except Exception:
        pass

    raise RuntimeError("Cannot parse PPT. Please save as PPTX and retry.")


# ── Other formats ──────────────────────────────────────────────

def parse_rtf(path):
    """RTF: striprtf library."""
    from striprtf.striprtf import rtf_to_text

    with open(path, "r", encoding="utf-8", errors="ignore") as f:
        rtf_content = f.read()
    text = rtf_to_text(rtf_content)
    title = os.path.splitext(os.path.basename(path))[0] or "RTF"
    return {"sourcePath": path, "docTitle": title, "sections": [
        {"heading": "", "level": 0, "body": text.strip(), "page": None}
    ]}


def parse_wps(path):
    """WPS: try as DOCX first (newer WPS = OOXML), fall back to DOC handler."""
    try:
        return parse_docx(path)
    except Exception:
        return parse_doc(path)


def parse_et(path):
    """ET: try as XLSX first, fall back to XLS handler."""
    try:
        return parse_xlsx(path)
    except Exception:
        return parse_xls(path)


def parse_dps(path):
    """DPS: try as PPTX first, fall back to PPT handler."""
    try:
        return parse_pptx(path)
    except Exception:
        return parse_ppt(path)


# ── Dispatch ───────────────────────────────────────────────────

HANDLERS = {
    "docx": parse_docx,
    "xlsx": parse_xlsx,
    "pptx": parse_pptx,
    "doc": parse_doc,
    "xls": parse_xls,
    "ppt": parse_ppt,
    "rtf": parse_rtf,
    "wps": parse_wps,
    "et": parse_et,
    "dps": parse_dps,
}


def handle(req):
    rid = req.get("id")
    method = req.get("method")
    if method != "parse":
        return {"id": rid, "ok": False, "error": f"unknown method: {method}"}
    params = req.get("params", {}) or {}
    path = params.get("path", "")
    kind = params.get("kind", "")
    handler = HANDLERS.get(kind)
    if handler is None:
        return {"id": rid, "ok": False, "error": f"unsupported kind: {kind}"}
    if not os.path.exists(path):
        return {"id": rid, "ok": False, "error": f"file not found: {path}"}
    try:
        result = handler(path)
        return {"id": rid, "ok": True, "result": result}
    except Exception as e:  # noqa: BLE001
        return {"id": rid, "ok": False, "error": f"{type(e).__name__}: {e}"}


def main():
    for line in sys.stdin:
        line = line.strip()
        if not line:
            continue
        try:
            req = json.loads(line)
        except Exception as e:  # noqa: BLE001
            sys.stdout.write(
                json.dumps({"id": None, "ok": False, "error": f"bad json: {e}"},
                           ensure_ascii=False) + "\n")
            sys.stdout.flush()
            continue
        resp = handle(req)
        sys.stdout.write(json.dumps(resp, ensure_ascii=False) + "\n")
        sys.stdout.flush()


if __name__ == "__main__":
    main()
